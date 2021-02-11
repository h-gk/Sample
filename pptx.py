#!/usr/bin/env python
# coding: utf-8

# In[1]:


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Mm
from pptx.enum.text import MSO_AUTO_SIZE

###1枚目のスライドを作成###
prs = Presentation()

#スライドサイズを決定
slide_height = 190.5
slide_width = 338.7
prs.slide_height = Mm(slide_height)
prs.slide_width = Mm(slide_width)


###スライドを追加する関数定義###
def slide_insert():
    title_slide_layout = prs.slide_layouts[6]
    global slide
    slide = prs.slides.add_slide(title_slide_layout)

def title_text(a):
    #タイトルのテキストボックスを張り付ける
    textbox_title1 = a

    #タイトルのテキストボックスの位置、幅と高さを決める
    textbox_left = Mm(3)
    textbox_top = Mm(5)
    textbox_width = Mm(slide_width - 50)
    textbox_height = Mm(18)
    textbox_fontsize = 40

    #テキストボックスを配置する
    textbox = slide.shapes.add_textbox(textbox_left,
                                   textbox_top,
                                   textbox_width,
                                   textbox_height)

    #テキストボックスに書き込む
    textbox.text = textbox_title1

    #フォントサイズを適用
    textbox.text_frame.paragraphs[0].font.size = Pt(textbox_fontsize)
    textbox.text_frame.paragraphs[0].font.bold = True


def msg_text(a):
    #タイトルのテキストボックスを張り付ける
    textbox_title2 = a

    #タイトルのテキストボックスの位置、幅と高さを決める
    textbox_left = Mm(3)
    textbox_top = Mm(23)
    textbox_width = Mm(slide_width - 20)
    textbox_height = Mm(13)
    textbox_fontsize = 20

    #テキストボックスを配置する
    textbox = slide.shapes.add_textbox(textbox_left,
                                   textbox_top,
                                   textbox_width,
                                   textbox_height)

    #テキストボックスに書き込む
    textbox.text = textbox_title2

    #フォントサイズを適用
    textbox.text_frame.paragraphs[0].font.size = Pt(textbox_fontsize)
    textbox.text_frame.word_wrap = True


def fig_text(a):
    #タイトルのテキストボックスを張り付ける
    textbox_title2 = a

    #タイトルのテキストボックスの位置、幅と高さを決める
    textbox_left = Mm(13)
    textbox_top = Mm(100)
    textbox_width = Mm(slide_width - 20)
    textbox_height = Mm(15)
    textbox_fontsize = 25

    #テキストボックスを配置する
    textbox = slide.shapes.add_textbox(textbox_left,
                                   textbox_top,
                                   textbox_width,
                                   textbox_height)

    #テキストボックスに書き込む
    textbox.text = textbox_title2

    #フォントサイズを適用
    textbox.text_frame.paragraphs[0].font.size = Pt(textbox_fontsize)
    textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def refe_text(a):
    #タイトルのテキストボックスを張り付ける
    textbox_title5 = '出所: ' + a 

    #タイトルのテキストボックスの位置、幅と高さを決める
    textbox_left = Mm(3)
    textbox_top = Mm(slide_height - 10)
    textbox_width = Mm(slide_width - 50)
    textbox_height = Mm(5)
    textbox_fontsize = 12

    #テキストボックスを配置する
    textbox = slide.shapes.add_textbox(textbox_left,
                                   textbox_top,
                                   textbox_width,
                                   textbox_height)

    #テキストボックスに書き込む
    textbox.text = textbox_title5

    #フォントサイズを適用
    textbox.text_frame.paragraphs[0].font.size = Pt(textbox_fontsize)



def cre_text(a):
    #タイトルのテキストボックスを張り付ける
    textbox_title6 = 'copy all rights reserved ● ●' + a

    #タイトルのテキストボックスの位置、幅と高さを決める
    textbox_left = Mm(slide_width -80)
    textbox_top = Mm(0)
    textbox_width = Mm(50)
    textbox_height = Mm(5)
    textbox_fontsize = 12

    #テキストボックスを配置する
    textbox = slide.shapes.add_textbox(textbox_left,
                                   textbox_top,
                                   textbox_width,
                                   textbox_height)

    #テキストボックスに書き込む
    textbox.text = textbox_title6

    #フォントサイズを適用
    textbox.text_frame.paragraphs[0].font.size = Pt(textbox_fontsize)

#--------プログラム展開部分-----------

# ファイルの開閉
f = open('slide.txt')
data = f.read()
f.close()

# 行数の確認
number = data.count('t:')
print(number)

#テキストファイルのキー文字を削除
data = data.replace('t:','')
data = data.replace('m:','')
data = data.replace('c:','')
data = data.replace('f:','').replace('r:','')


# 改行で区切る
lines = data.split('\n')
text = lines

for i in range(1,number + 1):
    #スライド挿入
    slide_insert()

    #挿入ファイルの有無確認
    len_text = len(text[5*i-1])

    #挿入ファイルがある場合だけ挿入
    if len_text > 0:
        slide.shapes.add_picture(text[5*i-1], Mm(50), Mm(50), Mm(100))

    title_text(text[5*i - 4])
    msg_text(text[5*i - 3])
    fig_text(text[5*i - 2])
    refe_text(text[5*i - 0])

prs.save('test.pptx')

